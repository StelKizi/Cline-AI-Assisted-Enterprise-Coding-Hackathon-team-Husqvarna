import json
import re
import base64
import mimetypes
from pathlib import Path
from mcp.server.fastmcp import FastMCP

mcp = FastMCP("kyra")

# Resolve design_system/ relative to this file, not the process cwd.
# Claude Desktop (and uv run) can launch the process from any working directory,
# so __file__.parent is the only reliable anchor.
KYRA_ROOT = Path(__file__).resolve().parent
DS_PATH = KYRA_ROOT / "design_system"


def _resolve_allowed_path(user_path: str) -> Path:
    raw = user_path.strip().strip('"').strip("'")
    p = Path(raw)
    if not p.is_absolute():
        p = (KYRA_ROOT / p).resolve()
    else:
        p = p.resolve()

    # Hard safety boundary: only allow reading files inside kyra-mcp/ directory.
    # This avoids accidental access to arbitrary user files from the MCP tool.
    if KYRA_ROOT not in p.parents and p != KYRA_ROOT:
        raise ValueError(f"Access denied. Path must be within {str(KYRA_ROOT)}")
    return p


def _json_response(payload: dict) -> str:
    return json.dumps(payload, indent=2, ensure_ascii=False)


def _read_text_file(path: Path, max_chars: int) -> tuple[str, bool]:
    text = path.read_text(encoding="utf-8", errors="replace")
    if len(text) > max_chars:
        return text[:max_chars], True
    return text, False


def _read_binary_file_base64(path: Path, max_bytes: int) -> tuple[str, int, bool]:
    data = path.read_bytes()
    truncated = False
    if len(data) > max_bytes:
        data = data[:max_bytes]
        truncated = True
    return base64.b64encode(data).decode("ascii"), len(data), truncated


def _extract_pdf_text(path: Path, max_chars: int) -> tuple[str, bool]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "PDF support requires dependency 'pypdf'. "
            "Run: uv add pypdf"
        ) from e

    reader = PdfReader(str(path))
    chunks: list[str] = []
    for page in reader.pages:
        try:
            chunks.append(page.extract_text() or "")
        except Exception:
            chunks.append("")
        if sum(len(c) for c in chunks) >= max_chars:
            break

    text = "\n".join(chunks).strip()
    if len(text) > max_chars:
        return text[:max_chars], True
    return text, False


def _extract_pptx_text(path: Path, max_chars: int) -> tuple[str, bool]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise RuntimeError(
            "PowerPoint support requires dependency 'python-pptx'. "
            "Run: uv add python-pptx"
        ) from e

    pres = Presentation(str(path))
    chunks: list[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                chunks.append(text)
        if sum(len(c) for c in chunks) >= max_chars:
            break

    text = "\n".join(chunks).strip()
    if len(text) > max_chars:
        return text[:max_chars], True
    return text, False


# ---------------------------------------------------------------------------
# Data loading helpers
# ---------------------------------------------------------------------------

def load_tokens() -> dict:
    # Read fresh from disk on every call — no caching intentional,
    # so hot-edits to tokens.json are reflected without restarting the server.
    return json.loads((DS_PATH / "tokens.json").read_text())

def load_components() -> dict:
    return json.loads((DS_PATH / "components.json").read_text())


def _extract_first_opening_tag(code: str, component_name: str) -> str | None:
    needle = f"<{component_name}"
    idx = code.find(needle)
    if idx == -1:
        return None
    boundary = idx + len(needle)
    if boundary >= len(code):
        return None
    nxt = code[boundary]
    if nxt not in " \n\r\t>/":
        return None
    i = boundary
    in_quote: str | None = None
    start = idx
    while i < len(code):
        c = code[i]
        if in_quote:
            if c == "\\" and i + 1 < len(code):
                i += 2
                continue
            if c == in_quote:
                in_quote = None
            i += 1
            continue
        if c in "\"'":
            in_quote = c
            i += 1
            continue
        if c == "/" and i + 1 < len(code) and code[i + 1] == ">":
            return code[start : i + 2]
        if c == ">":
            return code[start : i + 1]
        i += 1
    return None


def _variant_from_open_tag(open_tag: str) -> str | None:
    m = re.search(r'\bvariant=["\']([^"\']+)["\']', open_tag)
    return m.group(1) if m else None


def _discover_components_in_jsx(code: str, registry: dict) -> list[str]:
    names = sorted(registry.keys(), key=len, reverse=True)
    if not names:
        return []
    pat = re.compile(r"<(" + "|".join(re.escape(n) for n in names) + r")\b")
    seen: set[str] = set()
    out: list[str] = []
    for m in pat.finditer(code):
        n = m.group(1)
        if n not in seen:
            seen.add(n)
            out.append(n)
    return out


def _single_score_dict(name: str, code: str, spec: dict) -> dict:
    """Structured score for one component name against full snippet (variant read from first <Name …> tag)."""
    checks: list[str] = []
    passed = 0
    total = 0

    total += 1
    open_tag = _extract_first_opening_tag(code, name)
    variant_used = _variant_from_open_tag(open_tag) if open_tag else None
    if not variant_used:
        checks.append(
            f'✗ Variant check — variant prop missing on <{name}>. Required. Allowed: {spec["allowed_variants"]}'
        )
    elif variant_used not in spec["allowed_variants"]:
        checks.append(f'✗ Variant check — "{variant_used}" not in {spec["allowed_variants"]}')
    else:
        checks.append(f'✓ Variant check — "{variant_used}" is valid for <{name}>')
        passed += 1

    total += 1
    violations = []
    for pattern in spec["forbidden_patterns"]:
        if pattern in code:
            violations.append(f'"{pattern}"')
    if violations:
        checks.append(
            f'✗ Token check ({name}) — forbidden patterns in snippet: {", ".join(violations)}. Use design tokens instead.'
        )
    else:
        checks.append(f"✓ Token check ({name}) — no forbidden patterns in snippet")
        passed += 1

    total += 1
    hex_colors = re.findall(r"#[0-9A-Fa-f]{3,6}", code)
    if hex_colors:
        checks.append(f"✗ Color token check — hardcoded colors found: {hex_colors}. Use color tokens.")
    else:
        checks.append("✓ Color token check — no hardcoded hex colors")
        passed += 1

    total += 1
    lower_code = code.lower()
    raw_tags = {
        "button": "<button",
        "input": "<input",
        "card": "<div",
        "link": "<a ",
        "heading": "<h1",
        "text": "<p",
        "select": "<select",
        "checkbox": 'type="checkbox"',
    }
    raw_tag = raw_tags.get(name.lower())
    if raw_tag and raw_tag in lower_code and f"<{name}" not in code:
        checks.append(f'✗ Structure check — raw HTML "{raw_tag}" used instead of <{name}> component')
    else:
        checks.append(f"✓ Structure check — <{name}> component used correctly")
        passed += 1

    status = "PASS" if passed == total else "FAIL"
    fixes: list[str] = []
    if status == "FAIL":
        failing = [c for c in checks if c.startswith("✗")]
        for f in failing:
            if "Variant" in f:
                fixes.append(f'On <{name}>: use variant="{spec["allowed_variants"][0]}"')
            if "Token check" in f or "Color token" in f:
                fixes.append(
                    f'Replace inline styles/colors with design tokens: {spec["allowed_color_tokens"][:2]}'
                )
            if "Structure" in f:
                fixes.append(f"Use <{name}> component, not raw HTML")
    fixes_unique: list[str] = []
    for x in fixes:
        if x not in fixes_unique:
            fixes_unique.append(x)
    fixes = fixes_unique

    return {
        "status": status,
        "passed": passed,
        "total": total,
        "checks": checks,
        "fixes": fixes,
    }


def _format_scorecard_text(name: str, spec: dict, r: dict) -> str:
    summary = f"Result: {r['status']} ({r['passed']}/{r['total']} checks)\n\n" + "\n".join(r["checks"])
    if r["status"] == "FAIL" and r["fixes"]:
        summary += "\n\nRequired fixes:\n" + "\n".join(f"→ {fix}" for fix in r["fixes"])
    return summary


# ---------------------------------------------------------------------------
# MCP tools — each function becomes a tool the AI agent can call.
# Docstrings are the tool descriptions surfaced to the agent; keep them clear.
# Return type is always str: MCP tool results are plain strings, not dicts.
# ---------------------------------------------------------------------------

@mcp.tool()
def get_tokens() -> str:
    """Get all design tokens: colors, spacing, typography, border radius."""
    tokens = load_tokens()
    return json.dumps(tokens, indent=2)


@mcp.tool()
def get_component_spec(component_name: str) -> str:
    """
    Get the full contract for a design system component.
    Returns allowed variants, sizes, props, color tokens, and usage examples.
    """
    components = load_components()
    name = component_name.strip()

    if name not in components:
        available = list(components.keys())
        return f"Component '{name}' not found. Available: {available}"

    spec = components[name]
    return json.dumps(spec, indent=2)


@mcp.tool()
def list_components() -> str:
    """List all available components in the design system."""
    components = load_components()
    result = []
    for name, spec in components.items():
        result.append(f"- {name}: {spec['description']}")
    return "\n".join(result)


@mcp.tool()
def audit_context(component_name: str, intent: str) -> str:
    """
    Get generation constraints before writing UI code.
    Call this BEFORE generating a component to get the rules the AI must follow.

    Args:
        component_name: Name of the component to generate (e.g. "Button")
        intent: What you're trying to do (e.g. "primary submit button for signup form")
    """
    components = load_components()
    name = component_name.strip()

    if name not in components:
        available = list(components.keys())
        return f"Component '{name}' not found. Available: {available}"

    spec = components[name]

    # Build a focused constraint envelope — only the fields the agent needs
    # before writing code. The full spec (get_component_spec) has more detail
    # but also more noise; this keeps the context window lean.
    rules = {
        "component": name,
        "intent": intent,
        "constraints": {
            "required_props": spec["required_props"],
            "allowed_variants": spec["allowed_variants"],
            "allowed_sizes": spec["allowed_sizes"],
            "allowed_color_tokens": spec["allowed_color_tokens"],
            "forbidden_patterns": spec["forbidden_patterns"],
            "allowed_props": spec["allowed_props"],
        },
        # Positive + negative examples help the agent pattern-match faster
        # than reading a list of rules.
        "correct_example": spec["examples"]["correct"],
        "incorrect_example": spec["examples"]["incorrect"],
        "instruction": (
            f"Use only allowed_variants and allowed_sizes. "
            f"Reference colors by token name only (e.g. color.brand.primary). "
            f"Never use hardcoded hex, rgb, or inline styles. "
            f"Always include required_props: {spec['required_props']}."
        ),
    }

    # JSON string, not a dict — MCP tool results must be serializable strings.
    return json.dumps(rules, indent=2)


@mcp.tool()
def run_compliance_scorecard(component_name: str, code: str) -> str:
    """
    Validate generated UI code against design system constraints.
    Call this AFTER generating a component to check if it's compliant.

    Args:
        component_name: Name of the component that was generated (e.g. "Button")
        code: The generated JSX/HTML code to validate
    """
    components = load_components()
    name = component_name.strip()

    if name not in components:
        available = list(components.keys())
        return f"Component '{name}' not found. Available: {available}"

    spec = components[name]
    r = _single_score_dict(name, code, spec)
    return _format_scorecard_text(name, spec, r)


@mcp.tool()
def run_compliance_bundle(code: str) -> str:
    """
    Validate JSX that contains multiple Kyra components.
    Discovers tags like <Button>, <Input>, … in document order and runs the same
    checks as run_compliance_scorecard once per component type (variant read from
    each component's first opening tag). Returns JSON with per-component results
    and an overall PASS/FAIL.
    """
    registry = load_components()
    order = _discover_components_in_jsx(code, registry)
    if not order:
        return json.dumps(
            {
                "mode": "bundle",
                "overall": "FAIL",
                "components": [],
                "items": [],
                "totalPassed": 0,
                "totalChecks": 0,
                "message": "No Kyra components found in code (expected tags like <Button>, <Input>, …).",
            },
            indent=2,
        )
    items: list[dict] = []
    total_passed = 0
    total_checks = 0
    for n in order:
        spec = registry[n]
        r = _single_score_dict(n, code, spec)
        total_passed += r["passed"]
        total_checks += r["total"]
        items.append({"component": n, "result": r})
    overall = "PASS" if all(i["result"]["status"] == "PASS" for i in items) else "FAIL"
    return json.dumps(
        {
            "mode": "bundle",
            "overall": overall,
            "components": order,
            "items": items,
            "totalPassed": total_passed,
            "totalChecks": total_checks,
        },
        indent=2,
    )


@mcp.tool()
def read_asset(asset_path: str, max_chars: int = 200_000, max_bytes: int = 5_000_000) -> str:
    """
    Read an asset file from the kyra-mcp directory.

    Supported:
    - Text: .md, .txt, .json
    - Images: .png, .jpg/.jpeg, .gif, .webp (returned as base64)
    - Documents: .pdf (text extracted), .pptx (text extracted)

    Notes:
    - For safety, paths are restricted to files under the kyra-mcp/ directory.
    - Large outputs are truncated (see max_chars/max_bytes).
    """
    try:
        path = _resolve_allowed_path(asset_path)
    except Exception as e:
        return _json_response({"ok": False, "error": str(e)})

    if not path.exists():
        return _json_response({"ok": False, "error": f"File not found: {str(path)}"})
    if not path.is_file():
        return _json_response({"ok": False, "error": f"Not a file: {str(path)}"})

    ext = path.suffix.lower()
    mime, _ = mimetypes.guess_type(str(path))
    mime = mime or "application/octet-stream"

    try:
        if ext in {".md", ".txt"}:
            content, truncated = _read_text_file(path, max_chars=max_chars)
            return _json_response(
                {
                    "ok": True,
                    "path": str(path),
                    "kind": "text",
                    "mime": mime,
                    "truncated": truncated,
                    "content": content,
                }
            )

        if ext == ".json":
            raw, truncated = _read_text_file(path, max_chars=max_chars)
            if truncated:
                return _json_response(
                    {
                        "ok": True,
                        "path": str(path),
                        "kind": "json",
                        "mime": "application/json",
                        "truncated": True,
                        "content": raw,
                    }
                )
            try:
                obj = json.loads(raw)
                return _json_response(
                    {
                        "ok": True,
                        "path": str(path),
                        "kind": "json",
                        "mime": "application/json",
                        "truncated": False,
                        "content": obj,
                    }
                )
            except Exception:
                return _json_response(
                    {
                        "ok": True,
                        "path": str(path),
                        "kind": "json",
                        "mime": "application/json",
                        "truncated": False,
                        "content": raw,
                        "warning": "Invalid JSON; returned raw text.",
                    }
                )

        if ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
            b64, used_bytes, truncated = _read_binary_file_base64(path, max_bytes=max_bytes)
            return _json_response(
                {
                    "ok": True,
                    "path": str(path),
                    "kind": "image",
                    "mime": mime,
                    "encoding": "base64",
                    "bytes_read": used_bytes,
                    "truncated": truncated,
                    "content_base64": b64,
                }
            )

        if ext == ".pdf":
            text, truncated = _extract_pdf_text(path, max_chars=max_chars)
            return _json_response(
                {
                    "ok": True,
                    "path": str(path),
                    "kind": "pdf_text",
                    "mime": "application/pdf",
                    "truncated": truncated,
                    "content": text,
                }
            )

        if ext == ".pptx":
            text, truncated = _extract_pptx_text(path, max_chars=max_chars)
            return _json_response(
                {
                    "ok": True,
                    "path": str(path),
                    "kind": "pptx_text",
                    "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "truncated": truncated,
                    "content": text,
                }
            )

        if ext == ".ppt":
            return _json_response(
                {
                    "ok": False,
                    "path": str(path),
                    "error": "Legacy .ppt is not supported. Please convert to .pptx.",
                }
            )

        return _json_response(
            {
                "ok": False,
                "path": str(path),
                "error": f"Unsupported file type: {ext}",
                "supported_extensions": [
                    ".md",
                    ".txt",
                    ".json",
                    ".png",
                    ".jpg",
                    ".jpeg",
                    ".gif",
                    ".webp",
                    ".pdf",
                    ".pptx",
                ],
            }
        )
    except Exception as e:
        return _json_response(
            {
                "ok": False,
                "path": str(path),
                "error": str(e),
            }
        )


def main():
    mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
